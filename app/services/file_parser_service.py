"""
PM Digital Employee - File Parser Service
项目经理数字员工系统 - 文件解析服务

v1.3.0新增：支持多种文档格式的文本提取
"""

import os
import tempfile
import json
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any, Dict, List, Optional, Callable

from app.core.logging import get_logger
from app.core.exceptions import ServiceError

logger = get_logger(__name__)


class UnsupportedFormatError(ServiceError):
    """不支持文件格式错误."""

    def __init__(self, extension: str):
        super().__init__(
            code="unsupported_format",
            message=f"不支持 {extension} 格式文件",
            details={"extension": extension},
        )


class FileParseError(ServiceError):
    """文件解析错误."""

    def __init__(self, message: str, details: Optional[Dict[str, Any]] = None):
        super().__init__(
            code="file_parse_error",
            message=message,
            details=details or {},
        )


class OCRServiceError(ServiceError):
    """OCR服务错误."""

    def __init__(self, message: str):
        super().__init__(
            code="ocr_service_error",
            message=message,
        )


@dataclass
class FileInfo:
    """文件信息."""

    path: str                    # 文件路径
    name: str                    # 文件名
    extension: str               # 扩展名（不含点）
    size: int                    # 文件大小（字节）
    content_type: Optional[str] = None  # MIME类型


@dataclass
class TableData:
    """表格数据."""

    headers: List[str]           # 表头
    rows: List[List[str]]        # 行数据
    title: Optional[str] = None  # 表格标题/描述


@dataclass
class ParsedContent:
    """解析后的内容."""

    text: str                                           # 文本内容
    tables: List[TableData] = field(default_factory=list)  # 表格列表
    structure: Dict[str, Any] = field(default_factory=dict)  # 结构信息
    metadata: Dict[str, Any] = field(default_factory=dict)   # 元数据
    has_images: bool = False                            # 是否包含图片
    ocr_confidence: Optional[float] = None              # OCR置信度（图片解析时）

    def to_dict(self) -> Dict[str, Any]:
        """转换为字典格式."""
        return {
            "text": self.text,
            "tables": [
                {
                    "headers": t.headers,
                    "rows": t.rows,
                    "title": t.title,
                }
                for t in self.tables
            ],
            "structure": self.structure,
            "metadata": self.metadata,
            "has_images": self.has_images,
            "ocr_confidence": self.ocr_confidence,
        }

    def get_text_summary(self, max_length: int = 2000) -> str:
        """获取文本摘要."""
        if len(self.text) <= max_length:
            return self.text
        return self.text[:max_length] + "..."

    def get_key_content(self) -> str:
        """提取关键内容用于分类."""
        # 包含表格信息的简化文本
        key_parts = []

        # 添加文本摘要
        key_parts.append(self.get_text_summary(1000))

        # 添加表格摘要
        for table in self.tables[:3]:  # 最多3个表格
            if table.headers:
                key_parts.append(f"表头: {', '.join(table.headers)}")
            if table.rows:
                key_parts.append(f"表格行数: {len(table.rows)}")

        return "\n".join(key_parts)


class FileParserService:
    """
    文件解析服务.

    支持多种文档格式的文本提取，包括：
    - 文本文档：DOCX、DOC、PDF、TXT、MD
    - 表格文档：XLSX、XLS、CSV
    - 演示文档：PPTX、PPT
    - 图片文档：JPG、JPEG、PNG、BMP（需OCR）
    """

    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS = {
        # 文本文档
        "docx", "doc", "pdf", "txt", "md",
        # 表格文档
        "xlsx", "xls", "csv",
        # 演示文档
        "pptx", "ppt",
        # 图片文档
        "jpg", "jpeg", "png", "bmp",
    }

    # 文件大小限制（50MB）
    MAX_FILE_SIZE = 50 * 1024 * 1024

    # 图片扩展名（需要OCR）
    IMAGE_EXTENSIONS = {"jpg", "jpeg", "png", "bmp"}

    def __init__(self) -> None:
        """初始化文件解析服务."""
        self._parsers: Dict[str, Callable] = {}
        self._register_parsers()

    def _register_parsers(self) -> None:
        """注册解析器."""
        # 文本文档解析器
        self._parsers["docx"] = self._parse_docx
        self._parsers["doc"] = self._parse_doc_legacy
        self._parsers["pdf"] = self._parse_pdf
        self._parsers["txt"] = self._parse_text
        self._parsers["md"] = self._parse_text

        # 表格文档解析器
        self._parsers["xlsx"] = self._parse_xlsx
        self._parsers["xls"] = self._parse_xls
        self._parsers["csv"] = self._parse_csv

        # 演示文档解析器
        self._parsers["pptx"] = self._parse_pptx
        self._parsers["ppt"] = self._parse_ppt_legacy

        # 图片文档解析器（OCR）
        for ext in self.IMAGE_EXTENSIONS:
            self._parsers[ext] = self._parse_image

    def is_supported(self, extension: str) -> bool:
        """检查文件格式是否支持."""
        return extension.lower() in self.SUPPORTED_EXTENSIONS

    def is_image(self, extension: str) -> bool:
        """检查是否为图片文件."""
        return extension.lower() in self.IMAGE_EXTENSIONS

    async def parse(self, file_info: FileInfo) -> ParsedContent:
        """
        解析文件内容.

        Args:
            file_info: 文件信息

        Returns:
            ParsedContent: 解析后的内容

        Raises:
            UnsupportedFormatError: 不支持的文件格式
            FileParseError: 文件解析失败
        """
        ext = file_info.extension.lower()

        # 格式检查
        if not self.is_supported(ext):
            raise UnsupportedFormatError(ext)

        # 大小检查
        if file_info.size > self.MAX_FILE_SIZE:
            raise FileParseError(
                f"文件大小超过限制（{self.MAX_FILE_SIZE // 1024 // 1024}MB）",
                {"size": file_info.size, "max_size": self.MAX_FILE_SIZE},
            )

        logger.info(f"Parsing file: {file_info.name} (type: {ext})")

        try:
            parser = self._parsers[ext]
            content = await parser(file_info)

            # 设置元数据
            content.metadata["file_name"] = file_info.name
            content.metadata["file_extension"] = ext
            content.metadata["file_size"] = file_info.size
            content.metadata["parse_time"] = datetime.now().isoformat()

            logger.info(
                f"File parsed successfully: {file_info.name}",
                extra={
                    "text_length": len(content.text),
                    "table_count": len(content.tables),
                },
            )

            return content

        except (UnsupportedFormatError, FileParseError, OCRServiceError):
            raise
        except Exception as e:
            logger.error(f"Failed to parse file {file_info.name}: {e}")
            raise FileParseError(f"文件解析失败: {str(e)}", {"error": str(e)})

    # ==================== 文本文档解析 ====================

    async def _parse_docx(self, file_info: FileInfo) -> ParsedContent:
        """解析DOCX文件."""
        try:
            from docx import Document
        except ImportError:
            raise FileParseError("python-docx 库未安装，无法解析DOCX文件")

        doc = Document(file_info.path)

        text_parts = []
        tables = []

        # 提取段落文本
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # 提取表格
        for table in doc.tables:
            table_data = self._extract_docx_table(table)
            tables.append(table_data)

        return ParsedContent(
            text="\n".join(text_parts),
            tables=tables,
            structure={
                "type": "docx",
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
            },
        )

    def _extract_docx_table(self, table: Any) -> TableData:
        """从DOCX表格提取数据."""
        rows_data = []
        headers = []

        for i, row in enumerate(table.rows):
            row_data = [cell.text.strip() for cell in row.cells]
            if i == 0:
                headers = row_data
            else:
                rows_data.append(row_data)

        return TableData(
            headers=headers,
            rows=rows_data,
        )

    async def _parse_doc_legacy(self, file_info: FileInfo) -> ParsedContent:
        """解析旧版DOC文件（需要转换工具）."""
        # DOC格式较难直接解析，建议用户转换为DOCX
        raise FileParseError(
            "旧版DOC格式支持有限，建议转换为DOCX格式",
            {"extension": "doc", "suggestion": "转换为DOCX格式"},
        )

    async def _parse_pdf(self, file_info: FileInfo) -> ParsedContent:
        """解析PDF文件."""
        try:
            import pdfplumber
        except ImportError:
            raise FileParseError("pdfplumber 库未安装，无法解析PDF文件")

        text_parts = []
        tables = []

        with pdfplumber.open(file_info.path) as pdf:
            for page in pdf.pages:
                # 提取文本
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

                # 提取表格
                page_tables = page.extract_tables()
                for table_data in page_tables:
                    if table_data and len(table_data) > 0:
                        headers = table_data[0] if table_data else []
                        rows = table_data[1:] if len(table_data) > 1 else []
                        tables.append(TableData(
                            headers=[str(h) if h else "" for h in headers],
                            rows=[[str(c) if c else "" for c in row] for row in rows],
                        ))

        return ParsedContent(
            text="\n".join(text_parts),
            tables=tables,
            structure={
                "type": "pdf",
                "pages": len(pdf.pages) if pdf else 0,
            },
        )

    async def _parse_text(self, file_info: FileInfo) -> ParsedContent:
        """解析纯文本文件（TXT、MD）."""
        try:
            with open(file_info.path, "r", encoding="utf-8") as f:
                text = f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_info.path, "r", encoding="gbk") as f:
                    text = f.read()
            except Exception:
                raise FileParseError("无法识别文件编码")

        return ParsedContent(
            text=text,
            tables=[],
            structure={
                "type": "text",
                "lines": len(text.splitlines()),
                "chars": len(text),
            },
        )

    # ==================== 表格文档解析 ====================

    async def _parse_xlsx(self, file_info: FileInfo) -> ParsedContent:
        """解析XLSX文件."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise FileParseError("openpyxl 库未安装，无法解析XLSX文件")

        wb = load_workbook(file_info.path, read_only=True, data_only=True)

        text_parts = []
        tables = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]

            # 提取表格数据
            rows_data = []
            headers = []

            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if i == 0:
                    headers = row_data
                else:
                    if any(cell for cell in row_data):  # 非空行
                        rows_data.append(row_data)

            if headers and rows_data:
                tables.append(TableData(
                    headers=headers,
                    rows=rows_data,
                    title=sheet_name,
                ))

                # 生成表格文本描述
                text_parts.append(f"【{sheet_name}】")
                text_parts.append(f"列: {', '.join(headers)}")
                for row in rows_data[:20]:  # 最多20行用于文本
                    text_parts.append(" | ".join(row))

        wb.close()

        return ParsedContent(
            text="\n".join(text_parts),
            tables=tables,
            structure={
                "type": "xlsx",
                "sheets": len(wb.sheetnames),
            },
        )

    async def _parse_xls(self, file_info: FileInfo) -> ParsedContent:
        """解析旧版XLS文件."""
        try:
            import xlrd
        except ImportError:
            raise FileParseError("xlrd 库未安装，无法解析XLS文件")

        wb = xlrd.open_workbook(file_info.path)

        text_parts = []
        tables = []

        for sheet in wb.sheets():
            rows_data = []
            headers = []

            for row_idx in range(sheet.nrows):
                row_data = [
                    str(sheet.cell_value(row_idx, col_idx))
                    for col_idx in range(sheet.ncols)
                ]
                if row_idx == 0:
                    headers = row_data
                else:
                    rows_data.append(row_data)

            if headers and rows_data:
                tables.append(TableData(
                    headers=headers,
                    rows=rows_data,
                    title=sheet.name,
                ))

        return ParsedContent(
            text="\n".join(text_parts),
            tables=tables,
            structure={
                "type": "xls",
                "sheets": wb.nsheets,
            },
        )

    async def _parse_csv(self, file_info: FileInfo) -> ParsedContent:
        """解析CSV文件."""
        import csv

        rows_data = []
        headers = []

        # 尝试检测编码
        encodings = ["utf-8", "gbk", "gb2312", "utf-8-sig"]
        text_content = None

        for encoding in encodings:
            try:
                with open(file_info.path, "r", encoding=encoding) as f:
                    reader = csv.reader(f)
                    for i, row in enumerate(reader):
                        if i == 0:
                            headers = row
                        else:
                            rows_data.append(row)
                break
            except UnicodeDecodeError:
                continue
            except Exception as e:
                raise FileParseError(f"CSV解析失败: {str(e)}")

        if headers and rows_data:
            tables.append(TableData(
                headers=headers,
                rows=rows_data,
            ))

        # 生成文本
        text_parts = [f"列: {', '.join(headers)}"]
        for row in rows_data[:50]:
            text_parts.append(" | ".join(row))

        return ParsedContent(
            text="\n".join(text_parts),
            tables=tables,
            structure={
                "type": "csv",
                "rows": len(rows_data),
            },
        )

    # ==================== 演示文档解析 ====================

    async def _parse_pptx(self, file_info: FileInfo) -> ParsedContent:
        """解析PPTX文件."""
        try:
            from pptx import Presentation
        except ImportError:
            raise FileParseError("python-pptx 库未安装，无法解析PPTX文件")

        prs = Presentation(file_info.path)

        text_parts = []
        tables = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text)

                # PPTX中的表格
                if shape.has_table:
                    table_data = self._extract_pptx_table(shape.table)
                    tables.append(table_data)

            if slide_text:
                text_parts.append(f"【幻灯片 {slide_idx + 1}】")
                text_parts.extend(slide_text)

        return ParsedContent(
            text="\n".join(text_parts),
            tables=tables,
            structure={
                "type": "pptx",
                "slides": len(prs.slides),
            },
        )

    def _extract_pptx_table(self, table: Any) -> TableData:
        """从PPTX表格提取数据."""
        rows_data = []
        headers = []

        for i, row in enumerate(table.rows):
            row_data = [cell.text.strip() for cell in row.cells]
            if i == 0:
                headers = row_data
            else:
                rows_data.append(row_data)

        return TableData(
            headers=headers,
            rows=rows_data,
        )

    async def _parse_ppt_legacy(self, file_info: FileInfo) -> ParsedContent:
        """解析旧版PPT文件（需要转换工具）."""
        raise FileParseError(
            "旧版PPT格式支持有限，建议转换为PPTX格式",
            {"extension": "ppt", "suggestion": "转换为PPTX格式"},
        )

    # ==================== 图片OCR解析 ====================

    async def _parse_image(self, file_info: FileInfo) -> ParsedContent:
        """解析图片（OCR识别）."""
        # 调用OCR服务
        ocr_result = await self._call_ocr_service(file_info)

        return ParsedContent(
            text=ocr_result.text,
            tables=ocr_result.tables,
            structure={
                "type": "image",
                "format": file_info.extension,
            },
            has_images=True,
            ocr_confidence=ocr_result.confidence,
        )

    async def _call_ocr_service(self, file_info: FileInfo) -> ParsedContent:
        """
        调用OCR服务进行图片识别.

        当前使用模拟实现，实际部署时需对接真实OCR服务。
        可选方案：
        1. 百度OCR API
        2. 腾讯OCR API
        3. Azure Computer Vision
        4. Tesseract本地部署
        """
        # TODO: 对接实际OCR服务
        # 当前返回提示信息
        logger.warning(f"OCR service not configured for file: {file_info.name}")

        # 尝试使用 pytesseract（如果已安装）
        try:
            import pytesseract
            from PIL import Image

            image = Image.open(file_info.path)
            text = pytesseract.image_to_string(image, lang="chi_sim+eng")

            return ParsedContent(
                text=text,
                tables=[],
                structure={"type": "ocr", "engine": "tesseract"},
                has_images=True,
                ocr_confidence=0.75,  # Tesseract默认置信度估计
            )

        except ImportError:
            raise OCRServiceError(
                "OCR服务未配置，无法识别图片内容。"
                "请安装 pytesseract 或配置外部OCR服务。"
            )
        except Exception as e:
            raise OCRServiceError(f"OCR识别失败: {str(e)}")


# 服务工厂
_file_parser_service: Optional[FileParserService] = None


def get_file_parser_service() -> FileParserService:
    """获取文件解析服务实例."""
    global _file_parser_service
    if _file_parser_service is None:
        _file_parser_service = FileParserService()
    return _file_parser_service